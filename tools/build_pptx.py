# -*- coding: utf-8 -*-
"""生成提交用演示 PPTX：Private RAG Agent（深色精密仪器风 + AMD 红）。"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "submissions", "track2-private-rag-agent", "presentation.pptx")
os.makedirs(os.path.dirname(OUT), exist_ok=True)

# ---------- 主题 ----------
BG = RGBColor(0x10, 0x12, 0x18)       # 深色背景
BG2 = RGBColor(0x17, 0x1A, 0x22)      # 卡片底色
PANEL = RGBColor(0x1E, 0x22, 0x2C)    # 面板
RED = RGBColor(0xF5, 0x43, 0x4D)      # AMD 红
FG = RGBColor(0xE8, 0xEA, 0xF0)       # 主文字
MUT = RGBColor(0x9A, 0xA0, 0xB0)      # 次要文字
GREEN = RGBColor(0x4A, 0xC9, 0x8A)
AMBER = RGBColor(0xF2, 0xB0, 0x49)
LINE = RGBColor(0x2A, 0x2F, 0x3A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

FONT = "Microsoft YaHei"   # 中英混合显示稳妥


def _solid(shape, color, line=None, lw=0.75):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def _textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text, size=16, color=FG, bold=False, first=False, space_after=6,
          align=PP_ALIGN.LEFT, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def add_bg(slide):
    _solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height), BG)


def add_header(slide, title, kicker=None, page=None):
    """页眉：左上标题 + 底部 AMD 红线 + 页码。"""
    if kicker:
        tf = _textbox(slide, 0.7, 0.42, 11.9, 0.32)
        _para(tf, kicker.upper(), size=11, color=RED, bold=True, first=True)
    tf = _textbox(slide, 0.7, 0.72, 11.9, 0.75)
    _para(tf, title, size=28, color=FG, bold=True, first=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.52),
                                 Inches(0.9), Pt(3))
    _solid(bar, RED)
    if page:
        tf = _textbox(slide, 12.3, 7.05, 0.8, 0.35)
        _para(tf, str(page), size=11, color=MUT, align=PP_ALIGN.RIGHT, first=True)


def add_bullets(slide, items, x=0.7, y=1.85, w=11.9, h=5.0, size=15, gap=8):
    tf = _textbox(slide, x, y, w, h)
    first = True
    for it in items:
        if isinstance(it, tuple) and len(it) == 2:
            text, sub = it
        elif isinstance(it, tuple):
            text, sub = it[0], None
        else:
            text, sub = it, None
        p = _para(tf, "▪  " + text, size=size, color=FG, first=first, space_after=gap)
        first = False
        if sub:
            _para(tf, "      " + sub, size=size - 2, color=MUT, space_after=gap)
    return tf


def add_card(slide, x, y, w, h, title=None, title_color=FG):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    _solid(sh, BG2, line=LINE)
    if title:
        tf = sh.text_frame
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.14)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = title_color
        r.font.name = FONT
    return sh


# ================================================================ Slide 1 · Title
s = prs.slides.add_slide(BLANK)
add_bg(s)
# 顶部红色细条
_para(_textbox(s, 0.7, 0.75, 11.9, 0.3), "AMD AI DevMaster 2026 · TRACK 2", size=13,
      color=RED, bold=True, first=True)
_para(_textbox(s, 0.7, 2.0, 11.9, 1.1), "Private RAG Agent", size=54, color=FG, bold=True, first=True)
_para(_textbox(s, 0.7, 3.05, 11.9, 0.5),
      "A fully offline, local private multi-agent RAG assistant — built for AMD Radeon / ROCm",
      size=20, color=MUT, first=True)
for i, line in enumerate([
    "Decompose → Parallel Research → Fact-Check → Synthesize",
    "Data never leaves the machine · Every sentence auditable",
    "llama.cpp HIP / Ollama ROCm / vLLM · Radeon Cloud ready",
]):
    _para(_textbox(s, 0.7, 4.1 + i * 0.55, 11.9, 0.4), "— " + line, size=15,
          color=FG if i % 2 == 0 else MUT, first=True)
# 底部红色块 + 版本
_para(_textbox(s, 0.7, 6.7, 11.9, 0.4), "AMD AI DevMaster 2026 · Track 2 · MENG Yuxuan",
      size=12, color=MUT, first=True)

# ================================================================ Slide 2 · Problem
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Private Data Must Not Cross the Network", kicker="Problem", page=2)
add_bullets(s, [
    ("Enterprise and personal data — contracts, specs, internal docs — is precisely the content that cannot go to cloud LLM APIs (regulatory, confidentiality, IP risk).",),
    ("Cloud RAG compounds it: even \"hosted private\" means third-party custody of the index and query logs.",),
    ("Local complex-answer quality must match the cloud: multi-hop reasoning, long-document retrieval, verifiable citations.",),
    ("Gap: off-the-shelf local chatbots do single-turn retrieval with no grounding guarantees.",),
])
# 左右对照卡片
add_card(s, 0.7, 4.6, 5.6, 2.0, "Cloud AI", RED)
tf = _textbox(s, 0.95, 5.25, 5.1, 1.2)
_para(tf, "✗ Data egress to third-party APIs", size=14, color=FG, first=True)
_para(tf, "✗ Index + query logs in someone else's custody", size=14, color=FG)
_para(tf, "✗ Opaque provenance on answers", size=14, color=FG)
add_card(s, 6.9, 4.6, 5.6, 2.0, "This Project", GREEN)
tf = _textbox(s, 7.15, 5.25, 5.1, 1.2)
_para(tf, "✓ 100% local inference, zero egress", size=14, color=FG, first=True)
_para(tf, "✓ Deterministic grounding + citation checks", size=14, color=FG)
_para(tf, "✓ AMD Radeon GPU / ROCm optimized", size=14, color=FG)

# ================================================================ Slide 3 · Solution
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "A Fully Offline Private Multi-Agent RAG Assistant", kicker="Solution", page=3)
# 功能丝带
ribbon = ["OFFLINE", "MULTI-AGENT", "TRUST-VERIFIED", "AMD ROCm"]
x = 0.7
for r in ribbon:
    sh = slide_sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85),
                                       Inches(2.75), Inches(0.55))
    sh.adjustments[0] = 0.5
    _solid(sh, PANEL, line=RED if r == "AMD ROCm" else LINE)
    tf = sh.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = r
    run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = RED if r == "AMD ROCm" else FG; run.font.name = FONT
    x += 2.95
add_bullets(s, [
    ("Ingest PDF / Word / Markdown / Excel / PPT into a local knowledge base (demo corpus: 6 docs · 15 chunks).",),
    ("Answer complex questions via a four-stage pipeline: Decompose → Parallel Research → Fact-Check → Synthesize.",),
    ("Hybrid retrieval + cross-encoder rerank retrieves evidence, not just similar text.",),
    ("Two trust mechanisms — groundedness scoring + citation verification — make every answer auditable.",),
    ("Runs entirely on one machine: zero cloud dependencies, zero data egress.",),
], y=2.75)
# 微型框图
boxes = [("Private Docs", 0.7), ("Local KB", 3.4), ("RAG Agent", 6.1), ("Verified Answer", 8.8)]
for label, bx in boxes:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(5.6),
                            Inches(2.2), Inches(0.75))
    sh.adjustments[0] = 0.18
    _solid(sh, PANEL, line=RED if label == "Verified Answer" else LINE)
    tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = FG; r.font.name = FONT
    if bx < 8.8:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(bx + 2.25), Inches(5.85),
                                Inches(0.5), Inches(0.28))
        _solid(ar, MUT)

# ================================================================ Slide 4 · Architecture
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Agent Architecture — One Loop, Four Layers", kicker="Architecture", page=4)
layers = [
    ("UI  — FastAPI + SSE streaming · native HTML/CSS/JS · dark precision theme", FG),
    ("Agent  — memory · planner · tools (rag_search/read_doc/list_docs) · main loop", FG),
    ("RAG  — parse → chunk(512,64) → hybrid retrieval → cross-encoder rerank", FG),
    ("Trust  — groundedness · citation verification · LLM-as-judge", RED),
    ("Inference  — Ollama (ROCm) · llama.cpp (HIP/GGUF) · vLLM (ROCm)", FG),
]
y = 1.9
for label, color in layers:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(y),
                            Inches(10.9), Inches(0.78))
    sh.adjustments[0] = 0.16
    _solid(sh, PANEL, line=color)
    tf = sh.text_frame; tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.size = Pt(15); r.font.bold = (color == RED)
    r.font.color.rgb = FG if color != RED else RED; r.font.name = FONT
    y += 0.95
_para(_textbox(s, 1.2, 6.75, 10.9, 0.4), "Backend swap is config-only (config.yaml: model.backend = ollama | llama_cpp | vllm).",
      size=12, color=MUT, first=True)

# ================================================================ Slide 5 · Multi-Agent
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Decompose → Parallel Research → Fact-Check → Synthesize", kicker="Multi-Agent Pipeline", page=5)
# 横向管线
steps = [("Q", 0.7, RED), ("Decompose", 1.55, PANEL), ("n Researchers", 3.75, PANEL),
         ("Fact-Check", 6.6, PANEL), ("Verified Answer", 9.35, GREEN)]
for label, bx, bg in steps:
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(1.95),
                            Inches(1.75), Inches(0.85))
    sh.adjustments[0] = 0.16
    _solid(sh, bg, line=RED if bg == RED else (GREEN if bg == GREEN else LINE))
    tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = FG; r.font.name = FONT
add_bullets(s, [
    ("Decomposer splits the question into n independent sub-questions — no cross-dependencies, so they run in parallel (ThreadPoolExecutor).",),
    ("n Researcher sub-agents each retrieve the KB and emit a sub-report with its own source set.",),
    ("Fact-Checker applies two-stage verification per sub-report: deterministic groundedness + LLM-as-judge.",),
    ("Synthesizer composes the final answer exclusively from verified content.",),
    ("Final answer re-verified end-to-end (groundedness + citation precision) before display.",),
], y=3.15, size=14)
add_card(s, 0.7, 5.9, 11.9, 1.35, "Why parallelism wins on AMD GPU", RED)
tf = _textbox(s, 0.95, 6.5, 11.4, 0.7)
_para(tf, "GPU inference queues: serializing n research steps multiplies wall-clock time — parallelizing saturates GPU utilization.", size=13, color=FG, first=True)
_para(tf, "n agents share ONE embedding + ONE reranker (module-level singletons): parallelism costs a single VRAM copy. Native retrieval calls are lock-serialized (hnswlib/torch are not thread-safe).", size=13, color=MUT)

# ================================================================ Slide 6 · Retrieval
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Retrieval That Finds Evidence, Not Just Similar Text", kicker="Hybrid Retrieval", page=6)
add_bullets(s, [
    ("Query rewriting — a light local LLM call produces up to 2 synonym / precision variants.",),
    ("Vector search — ChromaDB (hnswlib, cosine) with bge-m3 embeddings: semantic recall.",),
    ("BM25 keyword — CJK-aware tokenization (2-gram sliding window): exact recall for proper nouns / abbreviations.",),
    ("Reciprocal Rank Fusion (RRF) merges both rankings rank-stably.",),
    ("Cross-encoder rerank (bge-reranker-v2-m3) scores (query, doc) jointly — the highest-value precision upgrade in RAG.",),
    ("Every stage degrades gracefully; pure vector search is the fallback floor.",),
], y=1.9, size=14)
# 机制对比表
from pptx.util import Inches as In
tbl_shape = s.shapes.add_table(5, 3, In(0.7), In(5.0), In(11.9), In(2.0))
tbl = tbl_shape.table
tbl.columns[0].width = In(2.8); tbl.columns[1].width = In(5.3); tbl.columns[2].width = In(3.8)
rows = [
    ("Mechanism", "What it catches", "Weakness"),
    ("Vector (bge-m3, cosine)", "Semantic paraphrase", "Misses rare terms / acronyms"),
    ("BM25 keyword", "Proper nouns, IDs, exact phrases", "No semantics"),
    ("RRF fusion", "Union of both, rank-stable", "—"),
    ("Cross-encoder rerank", "Top-k precision on (Q,D)", "Compute cost (small pool)"),
]
for ri, rdata in enumerate(rows):
    for ci, cell_text in enumerate(rdata):
        cell = tbl.cell(ri, ci)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        cell.text_frame.word_wrap = True
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = cell_text
        r.font.size = Pt(12); r.font.name = FONT
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = RED
        else:
            r.font.color.rgb = FG

# ================================================================ Slide 7 · Trust
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Auditable Answers: Groundedness + Citation Verification", kicker="Trust — The Differentiator", page=7)
add_bullets(s, [
    ("Groundedness — each answer sentence is checked against retrieved source text via character 3-gram overlap, labeled Supported / Partial / Unsupported with a 0–1 score.",),
    ("Citation verification — every [来源:file] reference must actually exist in the knowledge base; fabricated sources are rejected.",),
    ("LLM-as-judge (multi-agent mode) — a second pass reviews each factual claim against the sources.",),
    ("Synthesizer only uses sub-reports that pass verification; unverified sections are disclosed, not hidden.",),
    ("UI renders a per-sentence verification panel with a trust score — the judge sees exactly which sentence rests on which evidence.",),
], y=1.9, size=14)
# 例句展示
add_card(s, 0.7, 5.5, 11.9, 1.7, "Live panel (as rendered in the UI)", MUT)
tf = _textbox(s, 0.95, 6.15, 11.4, 1.0)
_para(tf, "✓ SUPPORTED 79%  — “Retrieval fuses vector + BM25 via Reciprocal Rank Fusion.”   [来源:技术白皮书.md]", size=13, color=GREEN, first=True)
_para(tf, "△ PARTIAL 41%  — “This section relies on model common knowledge; evidence is limited.”", size=13, color=AMBER)
_para(tf, "✓ Citation check — “技术白皮书.md” exists in the knowledge base (6 docs · 15 chunks).", size=13, color=GREEN)

# ================================================================ Slide 8 · AMD
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Multi-Backend, Quantized, VRAM-Slim — Built for Radeon", kicker="AMD / ROCm Optimization", page=8)
add_bullets(s, [
    ("Three interchangeable inference backends: llama.cpp (HIP, GGUF full-GPU offload) · Ollama (ROCm) · vLLM (ROCm, Radeon Cloud).",),
    ("GGUF quantization: Q4_K_M (~4.6 GB VRAM, fastest) as the balanced default; Q8_0 (~8.2 GB) for precision.",),
    ("VRAM sharing — embedding + reranker are module-level singletons: n parallel agents load one copy globally.",),
    ("Lazy loading — retrieval/rerank models load on first use, keeping startup VRAM low.",),
    ("Verification & monitoring — rocminfo · rocm-smi · ollama ps (confirms 100% GPU) · benchmarks/bench_amd.py.",),
], y=1.9, size=14)
tbl_shape = s.shapes.add_table(3, 4, In(0.7), In(5.15), In(11.9), In(1.7))
tbl = tbl_shape.table
tbl.columns[0].width = In(2.6); tbl.columns[1].width = In(2.8); tbl.columns[2].width = In(3.0); tbl.columns[3].width = In(3.5)
rows = [
    ("Quantization", "VRAM", "Speed", "Recommended use"),
    ("Q8_0", "~8.2 GB", "Fast", "Precision-critical"),
    ("Q4_K_M", "~4.6 GB", "Fastest", "Balanced default"),
]
for ri, rdata in enumerate(rows):
    for ci, ct in enumerate(rdata):
        cell = tbl.cell(ri, ci)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = ct
        r.font.size = Pt(13); r.font.name = FONT
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = RED
        elif ri == 2 and ci == 0:
            r.font.bold = True; r.font.color.rgb = GREEN
        else:
            r.font.color.rgb = FG

# ================================================================ Slide 9 · Performance
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "Measured Baseline, Documented Headroom on AMD", kicker="Performance", page=9)
add_bullets(s, [
    ("Baseline (dev machine, Ollama qwen3:8b, CPU): ≈ 1.8 token/s generation.",),
    ("AMD target (RX 7900 / ROCm, GGUF Q4_K_M, full-layer offload): projected 5–10× throughput.",),
    ("Automated, reproducible benchmark — python benchmarks/bench_amd.py (auto-detects GPU, reports tokens/s).",),
    ("Per-phase metrics: retrieval latency (ms) · rerank latency (ms) · generation throughput (tokens/s) · multi-agent end-to-end wall clock.",),
    ("Concurrency benefit — with GPU as the bottleneck, wall-clock scales with n parallel researchers.",),
], y=1.9, size=14)
tbl_shape = s.shapes.add_table(5, 4, In(0.7), In(4.9), In(11.9), In(2.3))
tbl = tbl_shape.table
tbl.columns[0].width = In(4.3); tbl.columns[1].width = In(2.5); tbl.columns[2].width = In(2.6); tbl.columns[3].width = In(2.5)
rows = [
    ("Workload", "CPU baseline", "AMD ROCm (measured)", "Speedup"),
    ("LLM generation (tokens/s)", "≈ 1.8", "to be filled", "—"),
    ("Embedding batch (ms)", "—", "to be filled", "—"),
    ("Hybrid retrieval + rerank (ms)", "—", "to be filled", "—"),
    ("Multi-agent end-to-end (s)", "—", "to be filled", "—"),
]
for ri, rdata in enumerate(rows):
    for ci, ct in enumerate(rdata):
        cell = tbl.cell(ri, ci)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = ct
        r.font.size = Pt(12); r.font.name = FONT
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = RED
        elif ct == "to be filled":
            r.font.color.rgb = AMBER
        else:
            r.font.color.rgb = FG
_para(_textbox(s, 0.7, 7.15, 11.9, 0.3),
      "Benchmarks: benchmarks/bench_amd.py · hardware: AMD RX 7900 XTX / ROCm 6.x / Q4_K_M",
      size=10, color=MUT, first=True)

# ================================================================ Slide 10 · Deployment
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "One Codebase, Three Deployment Paths", kicker="Deployment", page=10)
paths = [
    ("LOCAL", "pip install -r requirements.txt; ollama pull qwen3:8b + bge-m3; python main.py ui", FG),
    ("DOCKER", "docker compose up -d --build; uncomment devices block for ROCm GPU passthrough", FG),
    ("RADEON CLOUD", "deploy_rc.sh — 6-step one-shot: GPU detect → vLLM(ROCm) → config rewrite → smoke test", RED),
]
y = 1.9
for title, body, col in paths:
    add_card(s, 0.7, y, 11.9, 1.15, title, col)
    tf = _textbox(s, 0.95, y + 0.6, 11.4, 0.5)
    _para(tf, body, size=13, color=FG, first=True)
    y += 1.4
_para(_textbox(s, 0.7, 6.35, 11.9, 0.4),
      "Backend swap is config-only (model.backend: ollama | llama_cpp | vllm) — no code changes between local and cloud.",
      size=13, color=MUT, first=True)
_para(_textbox(s, 0.7, 6.75, 11.9, 0.4),
      "vLLM tuning on RC: HIP_VISIBLE_DEVICES=0 · TRITON_ATTN backend · int8 KV-cache · --gpu-memory-utilization 0.90.",
      size=13, color=MUT, first=True)

# ================================================================ Slide 11 · Demo
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_header(s, "UI That Shows the Work, Not Just the Answer", kicker="Demo Highlights", page=11)
add_bullets(s, [
    ("Streaming SSE chat with typewriter effect and a live activity timeline — every tool call and verification event appears in real time.",),
    ("GPU monitoring panel — utilization, VRAM, temperature, power, clock from pyamdgpuinfo / rocm-smi (demo fallback when no AMD GPU).",),
    ("Source deep-dive drawer — click any citation to open the original document with query terms highlighted.",),
    ("Per-sentence verification panel — Supported / Partial / Unsupported tags + trust score.",),
    ("Three-column dark UI, native HTML/CSS/JS + FastAPI — no heavy widget framework.",),
], y=1.9, size=14)

# ================================================================ Slide 12 · Closing
s = prs.slides.add_slide(BLANK)
add_bg(s)
_para(_textbox(s, 0.7, 2.2, 11.9, 0.6), "LOCAL · PARALLEL · VERIFIABLE", size=16,
      color=RED, bold=True, first=True)
_para(_textbox(s, 0.7, 2.8, 11.9, 1.0), "Private RAG Done Right", size=44, color=FG,
      bold=True, first=True)
add_bullets(s, [
    ("Fully offline multi-agent RAG — decompose → parallel research → fact-check → synthesize, zero cloud dependency.",),
    ("Two independent trust mechanisms make answers auditable sentence by sentence.",),
    ("Purpose-built for AMD — multi-backend ROCm inference, quantization, shared-model VRAM, vLLM on Radeon Cloud.",),
    ("Reproducible benchmarks and a one-shot Radeon Cloud deployment script included.",),
], y=4.0, size=15)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.2), Inches(2.2), Pt(3))
_solid(bar, RED)
_para(_textbox(s, 0.7, 6.4, 11.9, 0.4), "Track 2 · MENG Yuxuan · Private RAG Agent — video + deck + repository",
      size=12, color=MUT, first=True)

prs.save(OUT)
print("PPTX_OK slides=", len(prs.slides._sldIdLst), "->", OUT)
