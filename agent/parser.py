# -*- coding: utf-8 -*-
"""文档解析：PDF/Word/PPT/Excel/Markdown/TXT -> 文本切片"""
import os
import re
from typing import List


def read_document(path: str) -> str:
    """按扩展名解析单个文档，返回全文文本。"""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if ext == ".docx":
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if ext in (".md", ".txt"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    if ext in (".csv", ".json"):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join(str(c) for c in row if c is not None))
        return "\n".join(parts)
    return ""


def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64) -> List[str]:
    """按字符切块（中文场景按字符切比按 token 简单可靠）。"""
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if len(text) <= chunk_size:
        return [text] if text else []
    chunks, start = [], 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        # 尽量在段落/句号处断开
        if end < len(text):
            cut = max(text.rfind("\n\n", start, end), text.rfind("。", start, end),
                      text.rfind(". ", start, end), text.rfind("\n", start, end))
            if cut > start + chunk_size // 2:
                end = cut + 1
        chunks.append(text[start:end].strip())
        # 已到文档尾部：取完剩余文本立即结束，
        # 否则 end=len(text) 时 start 每次只前进 1 字符，产生 ~overlap 个重复尾块
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return [c for c in chunks if c]
